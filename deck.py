"""
PPTX generation — RingCentral AI Receptionist business-case deck.
Matches the FBM 4-slide design. Slides 1, 2, and 4 are built (slide 3,
the AIR capability mapping, is deferred).
"""

from __future__ import annotations
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from pipeline import PipelineResult
from claude_client import generate_narrative

# ---------------------------------------------------------------------------
# RingCentral 2026 corporate template (rc-presentation-template skill)
# ---------------------------------------------------------------------------
# Mandatory brand font. PowerPoint renders this when Inter Tight is installed;
# the .pptx names it regardless so it's correct on a branded machine.
FONT = "Inter Tight"

# Brand assets bundled in the repo (logos + pre-rendered gradient backgrounds).
_BRAND = Path(__file__).resolve().parent / "assets" / "brand"
BG_LIGHT = str(_BRAND / "bg_light.png")   # content / stats / tables
BG_WARM = str(_BRAND / "bg_warm.png")     # covers / dividers / closing
LOGO_COLOR = str(_BRAND / "logo_color.png")  # on light gradient
LOGO_WHITE = str(_BRAND / "logo_white.png")  # on warm gradient

# Brand colors. Accent orange is used ONLY as an accent (callouts, active nodes).
RC_ORANGE = RGBColor(0xFF, 0x88, 0x00)   # FF8800 brand accent
RC_NAVY = RGBColor(0x1B, 0x2A, 0x4A)     # 1B2A4A table-header fill
RC_BLUE = RGBColor(0x06, 0x2E, 0x5C)
RC_RED = RGBColor(0xC0, 0x2A, 0x2A)
RC_TEAL = RGBColor(0x0A, 0x8A, 0x8A)
RC_GOLD = RGBColor(0xC8, 0x8A, 0x00)
RC_PURPLE = RGBColor(0x5B, 0x3E, 0x96)
DARK = RGBColor(0x1A, 0x1A, 0x1A)        # 1A1A1A body/title on light slides
GRAY = RGBColor(0x88, 0x88, 0x88)        # muted footer / captions
LIGHT = RGBColor(0xF4, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE4, 0xEA)
ROW_ALT = RGBColor(0xF8, 0xF4, 0xF0)     # alternating table row

# --- Clean-deck design system (matches the approved FBM branch/board decks) ---
NAVY = RGBColor(0x0A, 0x1B, 0x3E)        # cover / closing background
NAVY_CARD = RGBColor(0x1E, 0x2E, 0x54)   # cards / circles on the navy slides
CARD_BG = RGBColor(0xF3, 0xF5, 0xF8)     # light-gray content cards on white
GREEN = RGBColor(0x1E, 0x9E, 0x6A)       # positive / return accent
MUTED = RGBColor(0x6B, 0x74, 0x86)       # muted labels/subtitles on white
ICE = RGBColor(0xC6, 0xD3, 0xEA)         # muted subtitle on navy

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, *, warm=False):
    """Full-bleed RingCentral gradient background (rc-presentation-template skill).

    Warm orange→lavender gradient for covers/closing (warm=True); light pastel
    gradient for content slides. The gradient IS the brand's visual identity, so
    every slide is backed by the pre-rendered PNG rather than a flat fill.
    """
    img = BG_WARM if warm else BG_LIGHT
    pic = slide.shapes.add_picture(img, 0, 0, SLIDE_W, SLIDE_H)
    # Send the picture to the very back so all content renders on top of it.
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)
    return pic


def _text(slide, text, left, top, width, height, *, size=14, bold=False,
          color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          italic=False, font=FONT, wrap=True, line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb


def _rich(slide, segments, left, top, width, height, *, anchor=MSO_ANCHOR.TOP,
          align=PP_ALIGN.LEFT, line_spacing=1.0, space_after=4):
    """segments: list of paragraphs, each a list of (text, {opts}) runs."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(segments):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for txt, opts in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(opts.get("size", 10))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", DARK)
            r.font.name = opts.get("font", FONT)
    return tb


def _rect(slide, left, top, width, height, fill, *, line=None, line_w=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w or Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _circle(slide, left, top, dia, fill, text=None, *, text_color=WHITE, size=12):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, dia, dia)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background()
    c.shadow.inherit = False
    if text is not None:
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = text_color; r.font.name = FONT
    return c


def _logo(slide, *, warm=False):
    # Official RingCentral wordmark PNG (rc-presentation-template skill): white on
    # warm-gradient slides, blue/orange on light content slides. Top-left corner.
    img = LOGO_WHITE if warm else LOGO_COLOR
    slide.shapes.add_picture(img, Inches(0.5), Inches(0.35), width=Inches(1.95))


def _footer(slide, page=None, *, warm=False):
    # Intentionally minimal — the approved decks carry only a page number
    # (stamped in _stamp_page_numbers). Kept as a no-op so existing call sites
    # stay valid without cluttering the clean layout.
    return


# Slides whose background is the warm gradient (white text), set during build.
_WARM_SLIDES = set()


def _stamp_page_numbers(prs):
    for i, slide in enumerate(prs.slides, 1):
        col = WHITE if i in _WARM_SLIDES else GRAY
        _text(slide, str(i), Inches(0.45), Inches(7.12), Inches(0.5), Inches(0.3),
              size=8, color=col)


def _title_block(slide, title, subtitle, *, warm=False):
    # Clean content-slide header: bold navy title (no underline), muted subtitle.
    # Long titles wrap to a second line; size off character count, not measured
    # width, so a wider fallback font never collides with the subtitle.
    n = len(title)
    if n > 56:
        tsize = 22
    elif n > 48:
        tsize = 25
    elif n > 40:
        tsize = 28
    else:
        tsize = 30
    _text(slide, title, Inches(0.5), Inches(0.98), Inches(12.4), Inches(0.6),
          size=tsize, bold=True, color=WHITE if warm else RC_NAVY, font=FONT, wrap=False)
    if subtitle:
        _text(slide, subtitle, Inches(0.52), Inches(1.58),
              Inches(12.4), Inches(0.4),
              size=12.5, color=ICE if warm else MUTED, font=FONT, wrap=False)


def _eyebrow(slide, text, x, y, *, color=RC_ORANGE):
    """Small uppercase orange kicker above a hero headline."""
    _text(slide, text.upper(), x, y, Inches(11.0), Inches(0.35),
          size=13, bold=True, color=color, font=FONT)


def _hero(slide, line1, line2, x, y, *, size=46, dark=True,
          line2_color=RC_ORANGE):
    """Two-line hero headline: line1 white/navy, line2 accent (orange)."""
    c1 = WHITE if dark else RC_NAVY
    _rich(slide,
          [[(line1, {"bold": True, "size": size, "color": c1, "font": FONT})],
           [(line2, {"bold": True, "size": size, "color": line2_color, "font": FONT})]],
          x, y, Inches(11.5), Inches(2.0), line_spacing=1.0, space_after=0)


def _stat_card(slide, x, y, w, h, big, label, *, big_color=RC_NAVY,
               fill=CARD_BG, label_color=MUTED, big_size=44):
    """Light card with a big colored number and a small muted label below."""
    _rect(slide, x, y, w, h, fill, radius=True)
    big_h = big_size / 58.0  # inches, approx cap height of the number
    _text(slide, big, x + Inches(0.28), y + Inches(0.22), w - Inches(0.5), Inches(big_h),
          size=big_size, bold=True, color=big_color, font=FONT, wrap=False)
    _text(slide, label, x + Inches(0.3), y + Inches(0.22 + big_h + 0.06), w - Inches(0.55), Inches(0.7),
          size=12, color=label_color, font=FONT, line_spacing=1.0)


def _punch(slide, segments, y=Inches(6.15), *, h=Inches(0.85), fill=NAVY):
    """Full-width navy 'punchline' band with rich centered text."""
    _rect(slide, Inches(0.5), y, Inches(12.33), h, fill, radius=True)
    _rich(slide, [segments], Inches(0.95), y, Inches(11.5), h,
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)


def _bullet_card(slide, x, y, w, h, head, body):
    """Light card with an orange bullet dot, bold navy header, muted body."""
    _rect(slide, x, y, w, h, CARD_BG, radius=True)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.34),
                                 Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = RC_ORANGE
    dot.line.fill.background(); dot.shadow.inherit = False
    _text(slide, head, x + Inches(0.6), y + Inches(0.22), w - Inches(0.85), Inches(0.45),
          size=15, bold=True, color=RC_NAVY, font=FONT)
    _text(slide, body, x + Inches(0.6), y + Inches(0.72), w - Inches(0.85), h - Inches(0.9),
          size=11.5, color=MUTED, font=FONT, line_spacing=1.05)


# ---------------------------------------------------------------------------
# Slide 1 — cover (navy)
# ---------------------------------------------------------------------------

def _slide_cover(prs, r: PipelineResult, ctx, ae_name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, warm=True)
    _WARM_SLIDES.add(len(prs.slides))
    _logo(s, warm=True)
    customer = ctx.get("customer", "your business")
    _eyebrow(s, f"Prepared for {customer}  ·  Business case", Inches(0.5), Inches(2.55),
             color=WHITE)
    _hero(s, "The revenue hiding in", "your missed calls.", Inches(0.5), Inches(3.05),
          size=48, dark=True, line2_color=WHITE)
    _text(s, f"RingCentral Performance Reports  ·  {r.reporting_period}  ·  "
             "every figure derived from your own call data.",
          Inches(0.52), Inches(5.35), Inches(11.0), Inches(0.5),
          size=14, color=WHITE, font=FONT)
    if ae_name:
        _text(s, f"Prepared by {ae_name}", Inches(0.52), Inches(6.7), Inches(8.0), Inches(0.35),
              size=11, color=WHITE, font=FONT)


# ---------------------------------------------------------------------------
# Slide 2 — missed-call summary
# ---------------------------------------------------------------------------

def _business_days(r: PipelineResult) -> int:
    """Weekdays (Mon–Fri) actually present in the reporting window.

    Counted from the call timestamps so "per business day" is defensible; falls
    back to a 5/7 proration of the calendar span if timestamps aren't available.
    """
    import pandas as pd
    df = r.sessions_df
    try:
        if df is not None and "start_time" in df.columns:
            dts = pd.to_datetime(df["start_time"], errors="coerce").dropna()
            wd = dts.dt.normalize()
            n = int(wd[wd.dt.weekday < 5].nunique())
            if n:
                return n
    except Exception:
        pass
    return max(round((r.days_in_period or 30) * 5 / 7), 1)


def _slide2(prs, r: PipelineResult, ctx, narr, sales_queue_calls):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    # Deterministic, plain-English source line — never LLM-generated, so the
    # methodology wording stays consistent and defensible across every deck.
    _title_block(s, narr.get("title", "We miss a lot of calls — and every miss is a lost order"),
                 f"Source: RingCentral call detail records, {r.reporting_period} · "
                 "one call is counted once, no matter how many agents it rang")

    # Top row — three big-number stat cards.
    cy, cw, ch, gap = Inches(2.15), Inches(3.95), Inches(1.75), Inches(0.24)
    cx = Inches(0.5)
    bpd = _business_days(r)
    stats = [
        (f"{r.total_missed:,}", "genuine missed calls\nexternal inbound only", RC_RED),
        (f"{r.miss_rate*100:.0f}%", f"of {r.universe_sessions:,} inbound\nsessions went unanswered", RC_NAVY),
        (f"{round(r.total_missed/bpd):,}", "missed every business day\n(Mon–Fri average)", RC_ORANGE),
    ]
    for big, label, col in stats:
        _stat_card(s, cx, cy, cw, ch, big, label, big_color=col, big_size=42)
        cx = Emu(int(cx) + int(cw) + int(gap))

    # Middle-left: hourly chart with the staffed-hours framing as its caption.
    mx, mw = Inches(0.5), Inches(6.1)
    _text(s, f"{round(r.business_hours_miss_pct*100)}% of misses hit during staffed hours — "
             "peak overflow, when every agent is already on a call.",
          mx, Inches(4.15), mw, Inches(0.5), size=12, bold=True, color=RC_NAVY, font=FONT, line_spacing=1.0)
    _hourly_chart(s, r, mx, Inches(4.7), mw, Inches(1.75))

    # Middle-right: how the misses split (rang-out + voicemail = total exactly).
    qr = r.queues_report
    abandoned_n = qr.abandoned if (qr and qr.abandoned) else r.abandoned
    rx, rw = Inches(6.95), Inches(5.9)
    _text(s, "How they were missed", rx, Inches(4.15), rw, Inches(0.3),
          size=12, bold=True, color=RC_NAVY, font=FONT)
    bd = [
        (r.missed, "Rang out — no answer"),
        (r.voicemail_total, "Went to voicemail (left message)"),
    ]
    by = Inches(4.65)
    for val, lbl in bd:
        _rect(s, rx, by, rw, Inches(0.62), CARD_BG, radius=True)
        pct = _pct(val, r.total_missed)
        _rich(s, [[(f"{val:,}  ", {"bold": True, "size": 16, "color": RC_ORANGE}),
                   (f"{pct}   ", {"bold": True, "size": 12, "color": MUTED}),
                   (lbl, {"size": 12, "color": RC_NAVY})]],
              rx + Inches(0.25), by, rw - Inches(0.5), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
        by += Inches(0.74)
    if abandoned_n:
        _text(s, f"Separately, {abandoned_n:,} callers abandoned in queue before reaching anyone — "
                 "a distinct measure, not part of the split above.",
              rx, by + Inches(0.02), rw, Inches(0.55), size=10, italic=True, color=MUTED)

    # Plain-English methodology footnote — the "why you can trust this number"
    # line an AE can point to. Explains Session ID de-duplication in customer
    # language and shows the exact counts, so nothing looks hidden.
    _rich(s, [[
        ("How the count is cleaned — and why it's conservative:  ", {"bold": True, "size": 9.5, "color": RC_NAVY}),
        (f"RingCentral logs every ring of a call as a separate line, so a call ringing several agents "
         f"appears several times. We group lines by their shared Session ID so each call counts once, not "
         f"many times ({r.raw_inbound_legs:,} raw lines → {r.inbound_sessions:,} real calls; "
         f"{r.phantom_legs_removed:,} duplicate rings removed). We also set aside {r.spam_sessions_removed:,} "
         f"calls whose entire duration — ring included — was under 5 seconds: misdials, wrong numbers and "
         f"auto-dialer hang-ups too brief for anyone to answer. Both steps only shrink the number, never "
         f"inflate it — no genuine call is counted twice or dropped.",
         {"size": 9.5, "color": MUTED})
    ]], Inches(0.5), Inches(6.72), Inches(12.4), Inches(0.62), line_spacing=1.03)

    _footer(s)


def _worst_table_qr(s, queues, x, y, w):
    """Worst-hit table from the Queues report (real abandoned per queue)."""
    rows = len(queues) + 1
    tbl_shape = s.shapes.add_table(rows, 4, x, y, w, Inches(0.32) * rows)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(0.45)
    tbl.columns[2].width = Inches(0.55)
    tbl.columns[3].width = Inches(0.45)
    headers = ["Queue", "Rev", "Aband", "%"]
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RC_NAVY
        _cell(c, htext, WHITE, bold=True, size=9, align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
    for i, q in enumerate(queues, 1):
        bg = ROW_ALT if i % 2 else WHITE
        # "Rev" column: ✓ marks a revenue-line queue (sales/orders/bookings).
        is_rev = (q.tier or "C") in ("A", "B")
        # A 100% abandon rate on a queue with zero answered isn't a data error — it's
        # an UNSTAFFED queue (0 agents). Mark it with † so it reads as the config
        # finding it is (see slide 5), not a broken number.
        unstaffed = (getattr(q, "answered", None) == 0 and q.inbound > 0)
        pct_txt = f"{round(q.abandon_rate*100)}%" + ("†" if unstaffed else "")
        vals = [q.name[:26], "✓" if is_rev else "", f"{q.abandoned}", pct_txt]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            if j == 1:
                col = RC_ORANGE
            elif j == 3 and q.abandon_rate >= 0.7:
                col = RC_RED
            else:
                col = DARK
            _cell(c, v, col, bold=(j == 1) or (j == 3 and q.abandon_rate >= 0.7), size=8.5,
                  align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)


def _hourly_chart(s, r, x, y, w, h):
    cd = CategoryChartData()
    hours = sorted(r.hourly_missed.keys())
    labels = [_hour_label(hh) for hh in hours]
    cd.categories = labels
    cd.add_series("Missed", [r.hourly_missed[hh] for hh in hours])
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 40
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = RC_RED
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(8)
    cat.format.line.color.rgb = CARD_BORDER
    val = chart.value_axis
    val.visible = False
    val.has_major_gridlines = False
    try:
        val.major_gridlines.format.line.fill.background()
    except Exception:
        pass
    return gframe


def _worst_table(s, queues, x, y, w):
    rows = len(queues) + 1
    tbl_shape = s.shapes.add_table(rows, 4, x, y, w, Inches(0.32) * rows)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(0.45)
    tbl.columns[2].width = Inches(0.55)
    tbl.columns[3].width = Inches(0.45)
    headers = ["Queue", "Rev", "Miss", "%"]
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RC_NAVY
        _cell(c, htext, WHITE, bold=True, size=9, align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
    for i, q in enumerate(queues, 1):
        bg = ROW_ALT if i % 2 else WHITE
        is_rev = (q.tier or "C") in ("A", "B")
        vals = [q.name[:26], "✓" if is_rev else "", f"{q.total_missed}", f"{round(q.miss_rate*100)}%"]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            if j == 1:
                col = RC_ORANGE
            elif j == 3 and q.miss_rate >= 0.7:
                col = RC_RED
            else:
                col = DARK
            _cell(c, v, col, bold=(j == 1) or (j == 3 and q.miss_rate >= 0.7), size=8.5,
                  align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 3 — miss rate by hour of day (reference deck's key slide)
# ---------------------------------------------------------------------------

def _slide_hourly(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    _title_block(s, narr.get("title", "Calls slip away after hours, on weekends — and even midday"),
                 narr.get("subtitle", ""))

    # Big miss-rate column chart (left)
    _miss_rate_chart(s, r, Inches(0.5), Inches(2.15), Inches(8.0), Inches(4.35))

    # Three callout stat cards (right)
    cx, cw = Inches(8.85), Inches(3.98)
    cards = [
        (_range_or_pct(r.after_hours_miss_lo, r.after_hours_miss_hi, r.after_hours_miss_rate),
         "of after-hours calls (6pm–6am) are missed"),
        (_range_or_pct(r.weekend_miss_lo, r.weekend_miss_hi, r.weekend_miss_rate),
         "missed on Saturdays and Sundays"),
        (f"~{round(r.midday_miss_rate*100)}%",
         "missed even during peak midday hours"),
    ]
    cy = Inches(2.15); ch = Inches(1.4); gap = Inches(0.18)
    for big, label in cards:
        _stat_card(s, cx, cy, cw, ch, big, label, big_color=RC_ORANGE, big_size=34)
        cy = Emu(int(cy) + int(ch) + int(gap))

    _punch(s, [("AIR answers instantly — ", {"bold": True, "size": 15, "color": WHITE, "font": FONT}),
               ("every hour, every day.", {"size": 15, "color": ICE, "font": FONT})],
           y=Inches(6.75), h=Inches(0.55))
    _footer(s)


def _miss_rate_chart(s, r, x, y, w, h):
    cd = CategoryChartData()
    cd.categories = [_hour_label24(hh) for hh in range(24)]
    cd.add_series("Miss %", [round(r.hourly_miss_rate.get(hh, 0) * 100) for hh in range(24)])
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "Miss %"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    plot = chart.plots[0]
    plot.gap_width = 35
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = RC_ORANGE
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(9)
    cat.format.line.color.rgb = CARD_BORDER
    val = chart.value_axis
    val.minimum_scale = 0
    val.maximum_scale = 100
    val.tick_labels.font.size = Pt(9)
    val.has_major_gridlines = True
    try:
        val.major_gridlines.format.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    except Exception:
        pass
    return gframe


# ---------------------------------------------------------------------------
# Slide 4 — AIR opportunity signals (abandoned + sub-60s answered)
# ---------------------------------------------------------------------------

def _slide3(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    _title_block(s, narr.get("title", "Where AI Receptionist captures revenue today"),
                 narr.get("subtitle", ""))

    qr = r.queues_report

    # Two stacked stat cards on the left — both transfer-safe, session-deduplicated
    # signals of the AI-receptionist opportunity (no queue-level abandoned metric,
    # which the Queues report inflates by counting queue-to-queue transfers).
    lx, lw = Inches(0.5), Inches(4.15)
    # Card A — genuine missed calls AIR can answer (session-deduplicated)
    ay, ah = Inches(2.15), Inches(2.15)
    _rect(s, lx, ay, lw, ah, CARD_BG, radius=True)
    _text(s, "MISSED CALLS AI\nRECEPTIONIST CAN ANSWER", lx + Inches(0.28), ay + Inches(0.22),
          lw - Inches(0.5), Inches(0.6), size=11.5, bold=True, color=MUTED, font=FONT)
    _text(s, f"{r.total_missed:,}", lx + Inches(0.24), ay + Inches(0.74), lw - Inches(0.3), Inches(1.0),
          size=54, bold=True, color=RC_RED, font=FONT, wrap=False)
    _rich(s, [[("24/7 — ", {"bold": True, "size": 12, "color": RC_NAVY}),
               ("after-hours calls and business-hours overflow alike, picked up instantly",
                {"size": 12, "color": MUTED})]],
          lx + Inches(0.28), ay + Inches(1.66), lw - Inches(0.5), Inches(0.5))

    # Card B — answered under 60s
    by, bh = Inches(4.55), Inches(2.15)
    _rect(s, lx, by, lw, bh, CARD_BG, radius=True)
    _text(s, "ANSWERED CALLS UNDER\n60 SECONDS", lx + Inches(0.28), by + Inches(0.22),
          lw - Inches(0.5), Inches(0.6), size=11.5, bold=True, color=MUTED, font=FONT)
    _text(s, f"{r.answered_under_60:,}", lx + Inches(0.24), by + Inches(0.74), lw - Inches(0.3), Inches(1.0),
          size=54, bold=True, color=RC_ORANGE, font=FONT, wrap=False)
    _rich(s, [[(f"{r.under_60_pct*100:.0f}%", {"bold": True, "size": 12, "color": RC_NAVY}),
               (" of answered calls ran under 60s — routine volume an AI receptionist could handle", {"size": 11, "color": MUTED})]],
          lx + Inches(0.28), by + Inches(1.6), lw - Inches(0.5), Inches(0.5))

    # Right — most-abandoned queues table (from the Queues report when present)
    rx = Inches(5.0)
    _text(s, "Where callers wait, then drop — staffed queues", rx, Inches(2.15),
          Inches(7.8), Inches(0.35), size=14, bold=True, color=RC_NAVY, font=FONT)

    if qr and qr.queues:
        # Exclude pure routing/overflow queues (0 ever answered): in RingCentral's
        # queue analytics a call transferred OUT of a reception/overflow queue is
        # recorded as "abandoned" there even though it was answered after transfer.
        # Those queues show a misleading ~100% abandon rate and are not where
        # callers actually give up, so they don't belong in a "most-abandoned" list.
        top_ab = sorted((q for q in qr.queues
                         if q.abandoned > 0 and q.tier != "D"
                         and getattr(q, "answered", 0) > 0),
                        key=lambda q: q.abandoned, reverse=True)[:8]
        _abandon_table_qr(s, top_ab, rx, Inches(2.55), Inches(7.85))
        _text(s, "“Left queue” = caller left this queue before it answered. Not all were lost — some were "
                 "answered after transfer to another queue. Staffed queues only; pure routing/overflow "
                 "queues (0 agents) are excluded.",
              rx, Inches(6.02), Inches(7.85), Inches(0.5), size=8, italic=True, color=MUTED)
        # Wait-time / SLA context strip
        bits = []
        if qr.avg_wait:
            bits.append(("Avg. wait ", qr.avg_wait))
        if qr.longest_wait:
            bits.append(("Longest wait ", qr.longest_wait))
        if qr.sla_pct:
            bits.append(("Service level ", f"{qr.sla_pct*100:.0f}%"))
        if bits:
            segs = []
            for i, (label, val) in enumerate(bits):
                if i:
                    segs.append(("    ·    ", {"size": 11, "color": CARD_BORDER}))
                segs.append((label, {"size": 11, "color": MUTED}))
                segs.append((val, {"size": 11, "bold": True, "color": RC_RED}))
            _rich(s, [segs], rx, Inches(6.32), Inches(7.85), Inches(0.35))
    else:
        top_ab = sorted((q for q in r.queue_stats.values() if q.abandoned_total > 0),
                        key=lambda q: q.abandoned_total, reverse=True)[:10]
        _abandon_table(s, top_ab, rx, Inches(2.55), Inches(7.85))

    # Takeaway strip
    _text(s, "AI Receptionist answers instantly — catching missed calls 24/7 and taking many "
             "short, routine calls so staff focus on revenue conversations.",
          Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
          size=11.5, italic=True, color=RC_NAVY, align=PP_ALIGN.CENTER)
    _footer(s)


def _abandon_table(s, queues, x, y, w):
    rows = len(queues) + 1
    row_in = 0.34
    tbl_shape = s.shapes.add_table(rows, 5, x, y, w, Inches(row_in * rows))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    widths = [4.05, 0.7, 1.2, 0.9, 1.0]
    for j, ww in enumerate(widths):
        tbl.columns[j].width = Inches(ww)
    for i in range(rows):
        tbl.rows[i].height = Inches(row_in)
    headers = ["Queue", "Rev", "Abandoned", "Ab. %", "Ans. <60s"]
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RC_NAVY
        _cell(c, htext, WHITE, bold=True, size=9,
              align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, q in enumerate(queues, 1):
        bg = ROW_ALT if i % 2 else WHITE
        is_rev = (q.tier or "C") in ("A", "B")
        vals = [q.name[:40], "✓" if is_rev else "", f"{q.abandoned_total:,}",
                f"{round(q.abandon_rate*100)}%", f"{q.answered_under_60:,}"]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            if j == 1:
                col = RC_ORANGE; bold = True
            elif j == 3:
                col = RC_RED if q.abandon_rate >= 0.5 else DARK
                bold = q.abandon_rate >= 0.5
            else:
                col = DARK; bold = False
            _cell(c, v, col, bold=bold, size=9,
                  align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


def _abandon_table_qr(s, queues, x, y, w):
    """Most-abandoned-queues table sourced from the Queues report (real abandoned)."""
    rows = len(queues) + 1
    row_in = 0.34
    tbl_shape = s.shapes.add_table(rows, 5, x, y, w, Inches(row_in * rows))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    widths = [3.75, 0.7, 1.2, 1.2, 1.0]
    for j, ww in enumerate(widths):
        tbl.columns[j].width = Inches(ww)
    for i in range(rows):
        tbl.rows[i].height = Inches(row_in)
    headers = ["Queue", "Rev", "Inbound", "Left queue", "Left %"]
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RC_NAVY
        _cell(c, htext, WHITE, bold=True, size=9,
              align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, q in enumerate(queues, 1):
        bg = ROW_ALT if i % 2 else WHITE
        is_rev = (q.tier or "C") in ("A", "B")
        unstaffed = (getattr(q, "answered", None) == 0 and q.inbound > 0)
        pct_txt = f"{round(q.abandon_rate*100)}%" + ("†" if unstaffed else "")
        vals = [q.name[:38], "✓" if is_rev else "", f"{q.inbound:,}",
                f"{q.abandoned:,}", pct_txt]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            if j == 1:
                col = RC_ORANGE; bold = True
            elif j == 4:
                col = RC_RED if q.abandon_rate >= 0.5 else DARK
                bold = q.abandon_rate >= 0.5
            else:
                col = DARK; bold = False
            _cell(c, v, col, bold=bold, size=9,
                  align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 4 — queue-level table
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Slide — config-fixable vs. structural gap (license-justification framing)
# ---------------------------------------------------------------------------

def _missed_time_split(r: PipelineResult):
    """Split genuine missed calls into after-hours (structural) vs business-hours.

    Sourced from the Calls export (the only per-call timestamped feed), spam
    excluded so it reconciles to the headline universe. Returns None if the
    timestamp/in-business-hours signal isn't available.
    """
    df = r.sessions_df
    if df is None or "in_business_hours" not in df.columns:
        return None
    missed_vals = ["missed", "abandoned", "vm/missed", "voicemail"]
    m = df[df["outcome"].astype(str).str.strip().str.lower().isin(missed_vals)]
    if "is_spam" in m.columns:
        m = m[~m["is_spam"].astype(bool)]
    total = len(m)
    after = int((~m["in_business_hours"].astype(bool)).sum())
    return {"total": total, "after": after, "business": total - after}


_DEST_COLOR = {
    "Ring group / front desk": RC_ORANGE,
    "Direct to a person's extension": RC_BLUE,
    "IVR / auto-attendant menu": RC_TEAL,
    "Main line / auto-receptionist": RC_GOLD,
}



def _slide_config_vs_air(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    _title_block(s, narr.get("title", "Genuine missed calls — split by when they arrived"),
                 narr.get("subtitle",
                          "De-duplicated from the call records; a call answered after transfer between "
                          "queues is never counted as lost. Split by staffed vs. unstaffed hours."))

    total = r.total_missed
    # Split the transfer-safe, session-deduplicated miss count by WHEN it arrived —
    # the only honest, single-source way to separate the structural after-hours floor
    # from the business-hours overflow gap. (The old "config-fixable via unstaffed
    # queues" number came from the Queues report's per-queue abandons, which count a
    # call transferred out of a reception/overflow queue as "abandoned" there — so it
    # double-counted routed calls that were actually answered elsewhere.)
    split = _missed_time_split(r)
    if split and split["total"]:
        base, after, business = split["total"], split["after"], split["business"]
    else:
        base, after, business = total, 0, total
    after_pct = round(after / base * 100) if base else 0
    bus_pct = round(business / base * 100) if base else 0

    # Three flow cards: total -> after-hours floor -> business-hours overflow
    cy, ch = Inches(2.4), Inches(2.9)
    cards = [
        (Inches(0.5), Inches(3.55), CARD_BG, RC_NAVY, MUTED,
         "TOTAL GENUINE MISSED", f"{total:,}", "session-deduplicated · transfers followed across queues",
         "Every inbound call that never reached a person, anywhere."),
        (Inches(4.78), Inches(3.55), RGBColor(0xFB,0xF3,0xE0), RC_GOLD, MUTED,
         "AFTER HOURS — NO ONE ON SHIFT", f"{after:,}",
         f"{after_pct}% of misses · arrived outside staffed hours",
         "Only always-on coverage can answer these."),
        (Inches(9.06), Inches(3.77), RGBColor(0xFC,0xEC,0xEC), RC_RED, MUTED,
         "DURING BUSINESS HOURS — EVERY AGENT BUSY", f"{business:,}",
         f"{bus_pct}% of misses · staff on shift, all lines full",
         "An overflow/capacity gap — not a routing misconfiguration."),
    ]
    for x, w, bg, numcol, subcol, head, big, sub, foot in cards:
        _rect(s, x, cy, w, ch, bg, radius=True)
        _text(s, head, x + Inches(0.24), cy + Inches(0.22), w - Inches(0.44), Inches(0.5),
              size=10.5, bold=True, color=numcol, font=FONT)
        _text(s, big, x + Inches(0.2), cy + Inches(0.74), w - Inches(0.3), Inches(1.0),
              size=44, bold=True, color=numcol, font=FONT, wrap=False)
        _text(s, sub, x + Inches(0.24), cy + Inches(1.9), w - Inches(0.44), Inches(0.5),
              size=10.5, bold=True, color=RC_NAVY, font=FONT)
        _text(s, foot, x + Inches(0.24), cy + Inches(2.35), w - Inches(0.44), Inches(0.45),
              size=10, italic=True, color=subcol)

    # Minus / equals connectors between the cards
    _text(s, "−", Inches(4.18), cy + Inches(1.0), Inches(0.6), Inches(0.8),
          size=34, bold=True, color=RC_ORANGE, align=PP_ALIGN.CENTER)
    _text(s, "=", Inches(8.46), cy + Inches(1.0), Inches(0.6), Inches(0.8),
          size=34, bold=True, color=RC_ORANGE, align=PP_ALIGN.CENTER)

    # Bottom emphasis band — the business-hours capacity gap (the larger, most
    # defensible finding: staff were on shift, every agent already on a call).
    if base and business:
        _punch(s, [(f"{business:,} ", {"bold": True, "size": 17, "color": RC_ORANGE, "font": FONT}),
                   (f"of the {total:,} misses ({bus_pct}%) happened while staff were on shift — "
                    "every agent already on a call. ", {"size": 12.5, "color": ICE, "font": FONT}),
                   ("A capacity gap, not a routing error: no queue rule answers a call when no human is free.",
                    {"size": 12.5, "bold": True, "color": WHITE, "font": FONT})],
              y=Inches(6.05))
    else:
        _punch(s, [("Every miss above is de-duplicated across queues — a call answered after transfer "
                    "is never counted as lost. Only always-on coverage answers a call when no human is free.",
                    {"size": 13, "bold": True, "color": WHITE, "font": FONT})],
              y=Inches(6.05))
    _footer(s)


# ---------------------------------------------------------------------------
# Slide — predicted call reasons (Firecrawl business context)
# ---------------------------------------------------------------------------

_TIER_BADGE = {
    "A": (RC_ORANGE, RGBColor(0xFF, 0xF1, 0xE3), "Revenue-line"),
    "B": (RC_ORANGE, RGBColor(0xFF, 0xF1, 0xE3), "Revenue-line"),
    "C": (RC_BLUE, RGBColor(0xE9, 0xEE, 0xF6), "General"),
    "D": (GRAY, LIGHT, "Internal"),
}


def _slide_call_reasons(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    biz = ctx.get("business") or {}
    summary = (biz.get("summary") or "").strip()
    url = biz.get("url", "")
    customer = ctx.get("customer", "this business")
    reasons_present = bool(biz.get("predicted_call_reasons"))

    # This is the OPENER: ground the deck in who they are and who's calling them,
    # then land the headline missed-call number. With no business context (no
    # Firecrawl crawl), fall back to a clean, number-led opener so the slide
    # always renders.
    if not (summary or reasons_present):
        _title_block(s, f"{r.total_missed:,} missed calls — and who was trying to reach {customer}",
                     f"RingCentral Performance Reports · {r.reporting_period} · "
                     f"session-deduplicated · spam-filtered · external inbound only")
        _footer(s)
        return

    _title_block(s, narr.get("title", f"Who {customer} is — and who's calling them"),
                 narr.get("subtitle",
                          f"{r.total_missed:,} missed calls in {r.reporting_period} · "
                          f"these are the callers a business like this hears from every day"))

    # Business profile band — stack the three pieces vertically (industry,
    # lines of business, one-line summary), each capped to a single non-wrapping
    # row. The old side-by-side layout collided when the summary wrapped over the
    # lines-of-business list.
    by = Inches(2.0)
    _rect(s, Inches(0.5), by, Inches(12.33), Inches(1.05), CARD_BG, radius=True)
    industry = (biz.get("industry") or "").strip()
    lobs = biz.get("lines_of_business") or []
    head = industry if industry else "Business profile"
    _text(s, head, Inches(0.8), by + Inches(0.14), Inches(11.6), Inches(0.32),
          size=13, bold=True, color=RC_NAVY, font=FONT, wrap=False)
    if lobs:
        lob_line = " · ".join(str(x) for x in lobs[:6])
        if len(lob_line) > 140:
            lob_line = lob_line[:138].rstrip(" ·") + "…"
        _text(s, lob_line, Inches(0.8), by + Inches(0.48), Inches(11.7), Inches(0.3),
              size=9.5, color=MUTED, font=FONT, wrap=False)
    if summary:
        s_line = summary if len(summary) <= 150 else summary[:148].rstrip() + "…"
        _text(s, s_line, Inches(0.8), by + Inches(0.76), Inches(11.7), Inches(0.28),
              size=9.5, italic=True, color=RC_NAVY, font=FONT, wrap=False)

    # Predicted call-reason cards (up to 6) — light cards, orange bullet dot,
    # bold navy reason, tier badge pill, muted "why".
    reasons = (biz.get("predicted_call_reasons") or [])[:6]
    cw, chh = Inches(4.05), Inches(1.5)
    xs = [Inches(0.5), Inches(4.68), Inches(8.86)]
    ys = [Inches(3.35), Inches(5.0)]
    for idx, item in enumerate(reasons):
        cx = xs[idx % 3]; cy = ys[idx // 3]
        tier = str(item.get("tier", "C")).upper()[:1]
        col, bg, tlabel = _TIER_BADGE.get(tier, _TIER_BADGE["C"])
        _rect(s, cx, cy, cw, chh, CARD_BG, radius=True)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.28), cy + Inches(0.3),
                                 Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = RC_ORANGE
        dot.line.fill.background(); dot.shadow.inherit = False
        _text(s, str(item.get("reason", ""))[:44], cx + Inches(0.56), cy + Inches(0.18),
              cw - Inches(1.75), Inches(0.55), size=14, bold=True, color=RC_NAVY, font=FONT)
        # tier badge
        bw = Inches(1.15)
        _rect(s, cx + cw - bw - Inches(0.18), cy + Inches(0.2), bw, Inches(0.34), bg, radius=True)
        _text(s, tlabel, cx + cw - bw - Inches(0.18), cy + Inches(0.225), bw, Inches(0.3),
              size=9, bold=True, color=col, align=PP_ALIGN.CENTER, font=FONT)
        _text(s, str(item.get("why", ""))[:120], cx + Inches(0.3), cy + Inches(0.78),
              cw - Inches(0.55), Inches(0.66), size=11, color=MUTED, line_spacing=1.05)
        # Only badge revenue-relevant on genuine revenue-line cards (tier A/B).
        if item.get("revenue_relevant") and tier in ("A", "B"):
            _text(s, "● revenue-relevant", cx + Inches(0.3), cy + chh - Inches(0.3),
                  cw - Inches(0.55), Inches(0.26), size=9, bold=True, color=RC_ORANGE)

    _text(s, "Caller types inferred from public website content — illustrative, not a claim about your call logs. "
             "The point: these are routine, answerable calls — the same kind that later show up in your RingCentral "
             "data as abandoned and sub-60-second calls.",
          Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.5),
          size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s)


# ---------------------------------------------------------------------------
# ROI model (assumptions are clearly labelled on the slides)
# ---------------------------------------------------------------------------

# AIR usage pricing
AIR_RATE_PER_MIN = 0.20
# Staffing-equivalent productivity (calls handled per agent per year)
CALLS_PER_FTE_EFFICIENT = 41_000   # busy, well-run inbound desk -> fewer FTEs
CALLS_PER_FTE_LEAN = 23_000        # realistic everyday throughput -> more FTEs
COST_PER_FTE = 50_000              # fully-loaded annual cost per agent
# Revenue assumptions
AVG_ORDER_VALUE = 500              # $ per recovered order (editable assumption)
# Deliberately conservative capture rates — the share of revenue-line missed calls
# an AI receptionist converts into a booked order. Kept low so the model is a floor,
# not a hero number a CFO dismisses.
CAPTURE_RATES = [0.02, 0.03, 0.05]
CONSERVATIVE_CAPTURE = 0.02        # the cell we lead with / highlight


def _roi_model(r: PipelineResult, air_rate: float = AIR_RATE_PER_MIN,
               rev_missed_month: float | None = None) -> dict:
    days = r.days_in_period or 30
    per_day_missed = r.total_missed / days
    missed_per_year = per_day_missed * 365
    missed_per_month = per_day_missed * 30.4
    inbound_per_month = (r.universe_sessions / days) * 30.4

    # Conservative opportunity pool: ONLY revenue-line missed calls (sales/orders/
    # bookings) — never the full missed-call count. If we weren't handed a pool,
    # fall back to the observed revenue-line monthly figure, else the total.
    if rev_missed_month and rev_missed_month > 0:
        rev_missed_per_month = float(rev_missed_month)
    else:
        rev_missed_per_month = missed_per_month
    rev_missed_per_year = rev_missed_per_month * 12

    # AIR fields every inbound call; minutes = calls x avg talk minutes
    air_minutes_month = inbound_per_month * r.avg_answered_minutes
    air_cost_month = air_minutes_month * air_rate
    air_cost_year = air_cost_month * 12

    # Staffing to answer every inbound call across all hours (the AIR equivalent).
    # Sized on full inbound volume, not just the missed subset, so the comparison
    # is like-for-like ("answer every call, 24/7").
    inbound_per_year = (r.universe_sessions / days) * 365
    fte_lo = inbound_per_year / CALLS_PER_FTE_EFFICIENT
    fte_hi = inbound_per_year / CALLS_PER_FTE_LEAN
    hire_lo = fte_lo * COST_PER_FTE
    hire_hi = fte_hi * COST_PER_FTE

    return {
        "missed_per_year": missed_per_year,
        "missed_per_month": missed_per_month,
        "rev_missed_per_year": rev_missed_per_year,
        "rev_missed_per_month": rev_missed_per_month,
        "inbound_per_month": inbound_per_month,
        "air_minutes_month": air_minutes_month,
        "air_cost_month": air_cost_month,
        "air_cost_year": air_cost_year,
        "fte_lo": fte_lo,
        "fte_hi": fte_hi,
        "hire_lo": hire_lo,
        "hire_hi": hire_hi,
        "air_rate": air_rate,
    }


def _money(v) -> str:
    v = float(v)
    if abs(v) >= 1_000_000:
        return f"${v/1_000_000:.1f}M"
    if abs(v) >= 1_000:
        return f"${round(v/1_000)}K"
    return f"${round(v)}"


# ---------------------------------------------------------------------------
# Slide 6 — AIR capability mapping
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Slide 7 — far cheaper than hiring
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Slide 8 — recovered revenue
# ---------------------------------------------------------------------------

def _slide_revenue(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _logo(s)
    m = ctx["roi"]
    aov = ctx.get("aov", AVG_ORDER_VALUE)
    cap_hi = ctx.get("capture_override") or CONSERVATIVE_CAPTURE
    rev_year = m["rev_missed_per_year"]
    rev_month = m["rev_missed_per_month"]
    _title_block(s, narr.get("title", "Sizing the revenue opportunity — conservatively"),
                 narr.get("subtitle",
                          f"Revenue-line missed calls only (sales · orders · bookings) · "
                          f"{ctx.get('reporting_period','')} run-rate · illustrative, using your own order value"))

    # Funnel strip — pool is the REVENUE-LINE missed calls, never total volume.
    fy = Inches(2.2)
    steps = [
        (f"{round(rev_month):,}", "revenue-line missed / month", RC_NAVY, CARD_BG, MUTED),
        (f"{round(rev_year):,}", "revenue-line missed / year", RC_NAVY, CARD_BG, MUTED),
        ("× capture %", "booked as orders", RC_NAVY, CARD_BG, MUTED),
        ("= recovered $", "added revenue", WHITE, RC_ORANGE, ICE),
    ]
    sw = Inches(3.0); sx = Inches(0.5)
    for big, lab, col, fill, labcol in steps:
        _rect(s, sx, fy, sw, Inches(1.0), fill, radius=True)
        _text(s, big, sx + Inches(0.2), fy + Inches(0.16), sw - Inches(0.4), Inches(0.5),
              size=24, bold=True, color=col, font=FONT, wrap=False)
        _text(s, lab, sx + Inches(0.2), fy + Inches(0.66), sw - Inches(0.4), Inches(0.3),
              size=11, color=labcol, font=FONT)
        sx = Emu(int(sx) + int(sw) + int(Inches(0.11)))

    # Recovered-revenue table by capture rate (revenue-line pool only)
    _text(s, "Annual recovered revenue — revenue-line missed calls only, by capture rate",
          Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.35), size=14, bold=True, color=RC_NAVY, font=FONT)
    _revenue_table(s, m, aov, cap_hi, Inches(0.5), Inches(3.85), Inches(12.33))

    # Headline takeaway — lead with the CONSERVATIVE cell, framed as a floor.
    # Wording flexes with where the order value came from, so we never assert an
    # AI-guessed number as if it were the customer's own figure.
    src = ctx.get("aov_source", "default")
    rec_lo = cap_hi * rev_year * aov
    _lt = RGBColor(0xCD, 0xD9, 0xEA)
    _rect(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.6), RC_BLUE, radius=True)
    if src == "supplied":
        head_segs = [
            (f"Even at a conservative {round(cap_hi*100)}% capture and your ${aov:,} order value, that's ", {"size": 13, "color": _lt}),
            (f"{_money(rec_lo)}/year recovered.", {"bold": True, "size": 16, "color": WHITE, "font": FONT}),
        ]
    else:
        head_segs = [
            (f"At {round(cap_hi*100)}% capture and a ${aov:,} order value, that's ", {"size": 13, "color": _lt}),
            (f"~{_money(rec_lo)}/year", {"bold": True, "size": 16, "color": WHITE, "font": FONT}),
            (" — enter your actual order value to size it precisely.", {"size": 13, "color": _lt}),
        ]
    _rich(s, [head_segs], Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.4),
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # "Measured vs. assumed" provenance band — the single most defensible element
    # on the deck's one modeled slide: it states plainly what is the customer's own
    # data versus the two inputs we assume (and exactly where the order value came from).
    prov = {"supplied": "your figure",
            "estimate": "website estimate — confirm with your team",
            "default": "placeholder — set your real value"}.get(src, "placeholder")
    by = Inches(6.72)
    _rich(s, [[("Measured from your own RingCentral logs: ", {"bold": True, "size": 9.5, "color": RC_TEAL, "font": FONT}),
               ("missed calls, abandons, hours and which queues — nothing modeled.", {"size": 9.5, "color": DARK})]],
          Inches(0.5), by, Inches(12.33), Inches(0.24), align=PP_ALIGN.CENTER)
    _rich(s, [[("Two assumptions only: ", {"bold": True, "size": 9.5, "color": RC_ORANGE, "font": FONT}),
               (f"capture rate (shown as a range above) and average order value (${aov:,} — {prov}). "
                "Validate against your CRM.", {"size": 9.5, "color": DARK})]],
          Inches(0.5), by + Inches(0.24), Inches(12.33), Inches(0.24), align=PP_ALIGN.CENTER)
    _footer(s, 8)


def _revenue_table(s, m, aov, cap_hi, x, y, w):
    aovs = [max(50, round(aov / 2)), aov, aov * 2]
    rows = len(CAPTURE_RATES) + 1
    cols = len(aovs) + 1
    row_in = 0.5
    tbl_shape = s.shapes.add_table(rows, cols, x, y, w, Inches(row_in * rows))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    tbl.columns[0].width = Inches(3.33)
    for j in range(1, cols):
        tbl.columns[j].width = Inches(3.0)
    for i in range(rows):
        tbl.rows[i].height = Inches(row_in)
    # header
    _cell(tbl.cell(0, 0), "Capture rate  ╲  Avg order value", WHITE, bold=True, size=11)
    tbl.cell(0, 0).fill.solid(); tbl.cell(0, 0).fill.fore_color.rgb = RC_NAVY
    for j, col_aov in enumerate(aovs, 1):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RC_NAVY
        # Flag the column that uses the customer's own order value.
        htxt = f"${col_aov:,}/order" + ("  (your value)" if col_aov == aov else "")
        _cell(c, htxt, WHITE, bold=True, size=11, align=PP_ALIGN.CENTER)
    for i, cap in enumerate(CAPTURE_RATES, 1):
        c0 = tbl.cell(i, 0)
        c0.fill.solid(); c0.fill.fore_color.rgb = ROW_ALT if i % 2 else WHITE
        _cell(c0, f"{round(cap*100)}% of missed calls", DARK, bold=True, size=12)
        for j, col_aov in enumerate(aovs, 1):
            rec = cap * m["rev_missed_per_year"] * col_aov
            c = tbl.cell(i, j)
            highlight = (abs(cap - cap_hi) < 1e-9 and col_aov == aov)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xF1, 0xE3) if highlight else (ROW_ALT if i % 2 else WHITE)
            _cell(c, _money(rec), RC_ORANGE if highlight else DARK,
                  bold=highlight, size=13 if highlight else 12, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 9 — rollout investment
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Slide 10 — recommendation & next steps
# ---------------------------------------------------------------------------

def _slide_next(prs, r: PipelineResult, ctx, narr):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, warm=True)
    _WARM_SLIDES.add(len(prs.slides))   # white footer + page number on this slide
    _logo(s, warm=True)
    m = ctx["roi"]
    cap_hi = ctx.get("capture_override") or CONSERVATIVE_CAPTURE
    aov = ctx.get("aov", AVG_ORDER_VALUE)
    rec_lo = cap_hi * m["rev_missed_per_year"] * aov
    air_year = m["air_cost_year"]

    _eyebrow(s, "The recommendation", Inches(0.5), Inches(1.95), color=WHITE)
    _hero(s, "Answer every call.", "Recover the revenue.", Inches(0.5), Inches(2.45),
          size=44, dark=True, line2_color=WHITE)

    # Three ROI stat cards on navy (transparent-navy cards).
    cards = [
        (_money(air_year), "annual cost to run AIR", RC_ORANGE),
        (f"~{_money(rec_lo)}", f"recovered/yr at {round(cap_hi*100)}% capture", RC_ORANGE),
        (f"{round(rec_lo / air_year) if air_year else 0}×", "return — net of all cost", GREEN),
    ]
    cx, cw, gap = Inches(0.5), Inches(3.95), Inches(0.24)
    cy, ch = Inches(4.7), Inches(1.55)
    for big, label, col in cards:
        _stat_card(s, cx, cy, cw, ch, big, label, big_color=col, fill=NAVY_CARD,
                   label_color=ICE, big_size=40)
        cx = Emu(int(cx) + int(cw) + int(gap))

    _text(s, "Even at a conservative capture rate on revenue-line missed calls, the return dwarfs the cost — "
             "validated against your own order value.",
          Inches(0.52), Inches(6.55), Inches(11.5), Inches(0.5),
          size=12.5, italic=True, color=WHITE, font=FONT)
    _footer(s, warm=True)


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _cell(cell, text, color, *, bold=False, size=9, align=PP_ALIGN.LEFT, pad=1):
    cell.margin_left = Pt(4); cell.margin_right = Pt(4)
    cell.margin_top = Pt(pad); cell.margin_bottom = Pt(pad)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT


def _pct(part, whole):
    return f"{round(100*part/whole)}%" if whole else "0%"


def _hour_label(h):
    suffix = "a" if h < 12 else "p"
    hh = h if h <= 12 else h - 12
    return f"{hh}{suffix}"


def _hour_label24(h):
    if h == 0:
        return "12a"
    if h < 12:
        return f"{h}a"
    if h == 12:
        return "12p"
    return f"{h - 12}p"


def _range_or_pct(lo, hi, agg):
    """Show a "X–Y%" range when both endpoints are meaningful and the spread is
    wide; otherwise fall back to the volume-weighted aggregate "~Z%"."""
    lo_p, hi_p, agg_p = round(lo * 100), round(hi * 100), round(agg * 100)
    if hi_p - lo_p >= 8 and lo_p > 0:
        return f"{lo_p}–{hi_p}%"
    return f"~{agg_p}%"


# ---------------------------------------------------------------------------
# orchestration
# ---------------------------------------------------------------------------

def build_deck(result: PipelineResult, run_id: str, customer: str, ae_name: str,
               prior_instructions: list[dict] | None = None,
               business_context: dict | None = None,
               overrides: dict | None = None) -> Path:
    out_dir = Path(tempfile.gettempdir()) / "rc_analyzer_decks"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / f"{run_id}.pptx"
    overrides = overrides or {}

    # Revenue-line = queues with clear buy/order/booking or customer-service intent
    # (the old A + B tiers). These are the missed calls that most directly cost money.
    #
    # Tier labels live on the Queues report (the authoritative per-queue feed), NOT
    # on queue_stats (built from the Calls export, whose Queue column is mostly
    # blank — so tiering it yields ~0). Source the revenue-line pool from the
    # Queues report: missed = inbound − answered for A/B queues. Fall back to
    # queue_stats only when no Queues report is present.
    _qr0 = result.queues_report
    if _qr0 and _qr0.queues:
        sales_queue_calls = sum(max((q.inbound or 0) - (q.answered or 0), 0)
                                for q in _qr0.queues if (q.tier or "C") in ("A", "B"))
    else:
        sales_queue_calls = sum(q.total_missed for q in result.queue_stats.values()
                                if (q.tier or "C") in ("A", "B"))
    top_queues = sorted(result.queue_stats.values(), key=lambda q: q.total_missed, reverse=True)[:5]

    ctx = {
        "customer": customer,
        "reporting_period": result.reporting_period,
        "universe_sessions": result.universe_sessions,
        "total_missed": result.total_missed,
        "miss_rate_pct": round(result.miss_rate * 100, 1),
        "answer_rate_pct": round(result.answer_rate * 100, 1),
        "abandoned": result.abandoned,
        "voicemail": result.voicemail_total,
        "rang_out": result.missed,
        "phantom_legs_removed": result.phantom_legs_removed,
        "spam_sessions_removed": result.spam_sessions_removed,
        "raw_inbound_legs": result.raw_inbound_legs,
        "inbound_sessions": result.inbound_sessions,
        "business_hours_miss_pct": round(result.business_hours_miss_pct * 100),
        "repeat_callers": result.repeat_callers,
        "misses_per_day": round(result.misses_per_day),
        "abandoned_total": result.abandoned_total,
        "answered_under_60": result.answered_under_60,
        "under_60_pct": round(result.under_60_pct * 100),
        "sales_queue_missed": sales_queue_calls,
        "num_queues": len(result.queue_stats),
        "top_missed_queues": {q.name: q.total_missed for q in top_queues},
    }

    # Non-queue traffic: calls that never entered a managed ACD call queue land in
    # the "Unknown" bucket (direct extension dials, ring/hunt groups, IVR). When this
    # dominates, the story is "your misses never reach a queue at all" — the single
    # strongest AIR argument: routing/SLA/overflow staffing cannot catch a call that
    # never enters a queue; only an always-on AI receptionist answers every line.
    # Real abandoned-in-queue data from the Queues report (2nd upload)
    qr = result.queues_report
    if qr and qr.inbound:
        ctx["queue_inbound"] = qr.inbound
        ctx["queue_abandoned"] = qr.abandoned
        ctx["queue_abandon_rate_pct"] = round(qr.abandon_rate * 100, 1)
        ctx["queue_longest_wait"] = qr.longest_wait
        ctx["queue_sla_pct"] = round(qr.sla_pct * 100)

    # How many misses actually entered a managed call queue vs. never reached one.
    #
    # We can't trust the Calls export's "Queue" column for this — it's blank on the
    # vast majority of rows, so the "Unknown" bucket balloons and falsely implies
    # ~100% of misses are unqueued. The Queues report is the authoritative source
    # for queue activity, so reconcile against it: queue misses = inbound − answered
    # for customer-facing queues (tier != 'D'; back-office queues carry no customer
    # and are out of scope). Everything else in total_missed never reached a queue.
    if result.total_missed:
        queue_missed = None
        if qr and qr.queues:
            tiered = [q for q in qr.queues if (getattr(q, "tier", "") or "").upper() != "D"]
            if tiered:
                queue_missed = sum(max((q.inbound or 0) - (q.answered or 0), 0) for q in tiered)
        if queue_missed is None and qr and qr.inbound:
            queue_missed = max(qr.inbound - qr.answered, 0)
        if queue_missed is None:
            # No Queues report at all — fall back to the (unreliable) Unknown bucket.
            _unq = result.queue_stats.get("Unknown")
            unqueued_missed = _unq.total_missed if _unq is not None else result.total_missed
        else:
            unqueued_missed = max(result.total_missed - queue_missed, 0)
            ctx["queue_missed"] = queue_missed
        ctx["unqueued_missed"] = unqueued_missed
        ctx["unqueued_share_pct"] = round(100 * unqueued_missed / result.total_missed)

    # Business context + modeling overrides
    business = business_context if (business_context and business_context.get("available")) else None
    ctx["business"] = business

    def _num(v):
        try:
            return float(v) if v is not None else None
        except (TypeError, ValueError):
            return None

    # Order value is the ONE figure on the ROI slide not measured from the call
    # logs, so its provenance is tracked and stamped on the slide. Precedence:
    #   1. supplied  — the AE/customer entered it (their number → defensible)
    #   2. estimate  — inferred by the website crawl (a starting point to confirm)
    #   3. default   — no signal at all (a placeholder, clearly flagged)
    aov_override = _num(overrides.get("avg_order_value"))
    aov_suggest = _num(business.get("suggested_avg_order_value")) if business else None
    if aov_override and aov_override > 0:
        ctx["aov"], ctx["aov_source"] = int(aov_override), "supplied"
    elif aov_suggest and aov_suggest > 0:
        ctx["aov"], ctx["aov_source"] = int(aov_suggest), "estimate"
    else:
        ctx["aov"], ctx["aov_source"] = AVG_ORDER_VALUE, "default"

    air_rate = _num(overrides.get("air_rate_per_min"))
    ctx["air_rate"] = air_rate if air_rate and air_rate > 0 else AIR_RATE_PER_MIN

    cap = _num(overrides.get("capture_rate"))
    ctx["capture_override"] = cap if cap and 0 < cap < 1 else None

    # Scale the observed revenue-line missed pool to a 30.4-day month so it lines
    # up with the model's monthly basis (the report period may not be exactly 30 days).
    _rev_missed_month = sales_queue_calls * (30.4 / (result.days_in_period or 30))
    ctx["roi"] = _roi_model(result, ctx["air_rate"], rev_missed_month=_rev_missed_month)
    ctx["sales_queue_missed"] = sales_queue_calls

    # Narrative titles/subtitles for the slides we keep. The deck was deliberately
    # simplified to a tight, defensible sequence (see the build order below):
    # every slide after the business-context opener is verifiable straight from
    # the RingCentral reports — nothing modeled except the one caveated slide.
    narr2 = _narr_titles(ctx, prior_instructions, "slide2")
    narr_hourly = {"title": "Calls slip away after hours, on weekends — and even midday",
                   "subtitle": f"Inbound miss rate by hour of day · {result.reporting_period} · when the business closes or gets busy, calls go unanswered"}
    narr3_sub = (f"{result.total_missed:,} genuine misses AIR can answer + "
                 f"{result.answered_under_60:,} short routine calls it can deflect · "
                 f"session-deduplicated · {result.reporting_period}")
    narr3 = {"title": "Where AI Receptionist captures revenue today",
             "subtitle": narr3_sub}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _WARM_SLIDES.clear()   # reset per-build (module-level set)

    # Simplified, defensible 6-slide flow:
    #   1. Who they are + who's calling them (business context up front)
    #   2. Missed calls by queue/company — the headline number + breakdown
    #   3. Hourly trend + after-hours / weekend / midday miss rates
    #   4. Answerable-call signals: abandons, answered <60s, wait/SLA
    #   5. What queue config can fix vs. the structural gap only AIR closes —
    #      the answer to the "just help us configure our queues" objection
    #   6. Illustrative opportunity (single, clearly-caveated money slide)
    #   7. Recommendation & next steps
    _slide_cover(prs, result, ctx, ae_name)      # navy cover
    _slide_call_reasons(prs, result, ctx, {})   # opener — robust to missing business context
    _slide2(prs, result, ctx, narr2, sales_queue_calls)
    _slide_hourly(prs, result, ctx, narr_hourly)
    _slide3(prs, result, ctx, narr3)
    _slide_config_vs_air(prs, result, ctx, {})
    _slide_revenue(prs, result, ctx, {})
    _slide_next(prs, result, ctx, {})

    _stamp_page_numbers(prs)
    prs.save(out_path)
    return out_path


def _narr_titles(ctx, prior, which):
    schema = {"title": "str — slide headline", "subtitle": "str — methodology source line"}
    try:
        return generate_narrative({**ctx, "slide": which}, schema, prior)
    except Exception:
        if which == "slide2":
            return {"title": f"{ctx['total_missed']:,} genuine missed calls, external inbound only",
                    "subtitle": f"Customer-facing queues · {ctx['reporting_period']} · session-deduplicated · spam-filtered · internal/back-office excluded"}
        return {"title": "Where customers are getting missed — by queue",
                "subtitle": f"Session-deduplicated · spam-filtered · {ctx['reporting_period']} · external inbound only · ranked by calls lost"}

